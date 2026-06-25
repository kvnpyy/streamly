#!/usr/bin/env python3
"""Fill Samsung App UI Description template for Streamly (certification-grade)."""
from __future__ import annotations

import sys
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = Path(
    "/Users/kevinpayoyo/Downloads/App_Description_template_eng_v.1.42 2.pptx"
)
OUT = ROOT / "tv-apps/assets/samsung-icons/Streamly-App-UI-Description-v1.0.2.pptx"
SHOTS = ROOT / "tv-apps/assets/screenshots/samsung"

APP_NAME = "Streamly"
CP_NAME = "Streamly (iptvwebplayer.org)"
APP_VERSION = "1.0.2"
AUTHOR = "Kevin Payoyo"
TODAY = date.today().isoformat()

UI_STRUCTURE = """Whole UI Structure (all screens — navigation flow)

[Login Screen]
  ├─ Link with PIN (default on TV)
  ├─ Xtream tab
  └─ M3U URL tab
       ↓ (successful sign-in)
[Home Screen]
  ├─ → [Live TV Screen] → [Player Overlay] → Return → back
  ├─ → [Movies Screen] → [Movie Detail Screen] → [Player Overlay]
  ├─ → [Series Screen] → [Series Detail Screen] → [Episode list] → [Player Overlay]
  ├─ → [Search Screen] → results → detail/play
  ├─ → [Settings Screen]
  └─ → [Continue Watching] → resume → [Player Overlay]
[Player Overlay] — full-screen video; Return exits to previous screen"""

USAGE_SCENARIOS = """UC1 — Sign in with PIN (recommended on Samsung TV)
1. On [Login Screen], focus is on Link with PIN tab.
2. On phone/PC: open https://iptvwebplayer.org/login → sign in (UC2 or UC3).
3. On phone/PC: [Settings Screen] → Link a TV with a PIN → copy 6-digit code.
4. On TV [Login Screen]: enter code → Continue → [Home Screen] loads.

TEST ACCOUNT (required — login feature):
Server URL (all): https://iptvwebplayer.org/api/review-panel
Password (all): StreamlyReview2026

Parallel model-group testing — one username per TV (same playlist):
  MG1: samsung_review (or samsung_review_1)
  MG2: samsung_review_2
  MG3: samsung_review_3 … through samsung_review_12

No ads, premium, or account-specific content — all QA logins are identical
except username. Sample Live TV, Movies, Series (public test streams only).

PIN alternative: PC login with any account above → Settings → Link TV PIN.
App ID: 3202606046491

UC2 — Sign in with Xtream (direct on TV)
1. [Login Screen] → Xtream tab.
2. Enter Server URL, Username, Password.
3. Sign in → [Home Screen].

UC3 — Sign in with M3U URL
1. [Login Screen] → M3U URL tab.
2. Paste M3U playlist URL → Sign in → [Home Screen].

UC4 — Browse and play Live TV
1. [Home Screen] or top nav → [Live TV Screen].
2. Optional: select category filter pill.
3. Focus channel card → OK → [Player Overlay] plays live stream.
4. Return → exits player to [Live TV Screen].

UC5 — Browse and play Movies
1. Top nav → [Movies Screen].
2. Focus movie poster → OK → [Movie Detail Screen].
3. Focus Play → OK → [Player Overlay].
4. Return → back to detail or catalog.

UC6 — Browse and play Series
1. Top nav → [Series Screen].
2. Select series → [Series Detail Screen] → select episode → [Player Overlay].

UC7 — Search
1. Top nav → [Search Screen].
2. Enter query → select result → play or open detail.

UC8 — Settings / PIN pairing
1. Top nav → [Settings Screen].
2. Link a TV with a PIN (for phone-to-TV pairing).
3. Account, preferences, sign out.

UC9 — Exit player
1. While [Player Overlay] is visible, press Return → previous browse screen."""

MENU_LOGIN_SETTINGS = """[Login Screen] — screenshot 1
| # | Element | Action |
| 1 | Link with PIN tab | PIN entry (default on TV) |
| 2 | Xtream tab | Server / user / password fields |
| 3 | M3U URL tab | Playlist URL field |
| 4 | Continue | Pair TV → [Home Screen] |

[Settings Screen] — screenshot 2
| # | Element | Action |
| 1 | Link a TV with a PIN | Generate code on web for TV entry |
| 2 | Account row | View signed-in user |
| 3 | Playback quality | Adjust streaming preference |
| 4 | Sign out | End session on this TV |"""

MENU_HOME = """[Home Screen]
| # | Element | Action |
| 1 | Top nav links | Navigate to main screens |
| 2 | Continue watching tile | Resume or open detail |
| 3 | On now channel card | Tune live → [Player Overlay] |
| 4 | Featured shelf tile | Open detail or play |
| 5 | See all | Open full list for row |"""

MENU_LIVE_PLAYER = """[Live TV Screen] — screenshot 1
| # | Element | Action |
| 1 | Category filter pill | Filter channel list |
| 2 | Channel card | Play live → [Player Overlay] |
| 3 | Now playing panel | Program info (read-only) |

[Player Overlay] — screenshot 2
| # | Element | Action |
| 1 | Progress bar | Seek position (VOD) |
| 2 | Play / Pause | Toggle playback |
| 3 | Return (remote) | Exit player to previous screen |"""

MENU_CATALOG = """[Movies Screen] — top-left screenshot
| # | Element | Action |
| 1 | Genre filter | Filter movie grid |
| 2 | Poster card | Open [Movie Detail Screen] |

[Movie Detail Screen] — top-right
| # | Element | Action |
| 1 | Play | Start VOD → [Player Overlay] |
| 2 | My List | Toggle favorite |
| 3 | More like this | Open another title |

[Series Screen] — bottom-left
| # | Element | Action |
| 1 | Genre filter | Filter series grid |
| 2 | Series poster | Open series detail → episodes |

[Search Screen] — bottom-right
| # | Element | Action |
| 1 | Search field | Enter query |
| 2 | Result row | Play or open detail |"""


def set_text(shape, text: str, font_size: int = 14) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)


def set_cell(table, row: int, col: int, text: str) -> None:
    table.cell(row, col).text = text


def clear_placeholders(slide) -> None:
    to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            continue
        if hasattr(shape, "left"):
            to_remove.append(shape)
    for shape in to_remove:
        slide.shapes._spTree.remove(shape._element)


def add_pair(slide, left_img: Path, right_img: Path) -> None:
    if left_img.exists():
        slide.shapes.add_picture(
            str(left_img), Inches(0.35), Inches(2.55), width=Inches(4.55)
        )
    if right_img.exists():
        slide.shapes.add_picture(
            str(right_img), Inches(5.05), Inches(2.55), width=Inches(4.55)
        )


def add_quad(slide, imgs: list[Path]) -> None:
    positions = [
        (Inches(0.35), Inches(2.55), Inches(4.55)),
        (Inches(5.05), Inches(2.55), Inches(4.55)),
        (Inches(0.35), Inches(5.35), Inches(4.55)),
        (Inches(5.05), Inches(5.35), Inches(4.55)),
    ]
    for img, (left, top, width) in zip(imgs, positions):
        if img.exists():
            slide.shapes.add_picture(str(img), left, top, width=width)


def delete_slide(prs: Presentation, index: int) -> None:
    r_id = prs.slides._sldIdLst[index].rId
    prs.part.drop_rel(r_id)
    del prs.slides._sldIdLst[index]


def main() -> int:
    if not TEMPLATE.exists():
        print(f"Template not found: {TEMPLATE}", file=sys.stderr)
        return 1

    prs = Presentation(str(TEMPLATE))
    delete_slide(prs, 0)

    slide = prs.slides[0]
    set_text(slide.shapes[0], f"{APP_NAME} — Application UI Description")
    set_text(slide.shapes[1], CP_NAME)

    slide = prs.slides[1]
    tbl = slide.shapes[1].table
    set_cell(tbl, 1, 0, "1.0")
    set_cell(tbl, 1, 1, TODAY)
    set_cell(
        tbl,
        1,
        2,
        f"Initial Samsung TV certification submission\n"
        f"Application package version {APP_VERSION}",
    )
    set_cell(tbl, 1, 3, AUTHOR)

    slide = prs.slides[3]
    set_text(slide.shapes[1], UI_STRUCTURE, font_size=11)

    slide = prs.slides[4]
    set_text(slide.shapes[0], "Menu & function — [Login Screen] + [Settings Screen]")
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        set_text(slide.shapes[1], MENU_LOGIN_SETTINGS, font_size=8)
    clear_placeholders(slide)
    add_pair(slide, SHOTS / "01-login.jpg", SHOTS / "08-settings.jpg")

    slide = prs.slides[5]
    set_text(slide.shapes[0], "Menu & function — [Home Screen]")
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        set_text(slide.shapes[1], MENU_HOME, font_size=9)
    clear_placeholders(slide)
    slide.shapes.add_picture(
        str(SHOTS / "04-tv-home.jpg"), Inches(4.9), Inches(1.1), width=Inches(4.6)
    )

    slide = prs.slides[6]
    set_text(slide.shapes[0], "Menu & function — [Live TV Screen] + [Player Overlay]")
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        set_text(slide.shapes[1], MENU_LIVE_PLAYER, font_size=8)
    clear_placeholders(slide)
    add_pair(slide, SHOTS / "02-live-tv.jpg", SHOTS / "09-player.jpg")

    slide = prs.slides[7]
    set_text(slide.shapes[0], "Usage Scenario — Streamly")
    set_text(slide.shapes[1], USAGE_SCENARIOS, font_size=9)
    if len(slide.shapes) > 2:
        try:
            slide.shapes._spTree.remove(slide.shapes[2]._element)
        except Exception:
            pass

    slide = prs.slides[8]
    set_text(
        slide.shapes[0],
        "Menu & function — [Movies] [Movie Detail] [Series] [Search]",
    )
    set_text(slide.shapes[1], MENU_CATALOG, font_size=7)
    clear_placeholders(slide)
    add_quad(
        slide,
        [
            SHOTS / "05-movies-grid.jpg",
            SHOTS / "03-movie-detail.jpg",
            SHOTS / "06-series-grid.jpg",
            SHOTS / "07-search.jpg",
        ],
    )

    slide = prs.slides[9]
    set_text(slide.shapes[0], "Key Policy — Streamly")
    set_text(
        slide.shapes[1],
        "Standard Tizen keys. Return/Exit not modified. Volume uses TV system.",
        font_size=11,
    )
    tbl = slide.shapes[2].table
    policies = [
        ("UP / DOWN / LEFT / RIGHT", "Move focus between UI elements", ""),
        ("OK / ENTER", "Activate focused button; play content", ""),
        ("Return", "Back one screen; exit [Player Overlay]", "Samsung Mandatory"),
        ("Exit", "Close application (system)", "Samsung Mandatory"),
        ("Play / Pause", "Toggle playback in player when focused", ""),
        ("Ch. Up/Down", "Not used by app (N/R)", ""),
        ("Color / Info / Tool keys", "Not used by app (N/R)", ""),
        ("Volume +/-", "TV system volume — not overridden by app", ""),
    ]
    for i, (btn, action, remark) in enumerate(policies):
        r = i + 2
        if r < len(tbl.rows):
            set_cell(tbl, r, 0, btn)
            set_cell(tbl, r, 1, action)
            set_cell(tbl, r, 2, remark)

    slide = prs.slides[10]
    tbl = slide.shapes[1].table
    set_cell(
        tbl,
        1,
        1,
        "The application has no language options. English UI only. "
        "No in-app language selector. Content metadata language depends on the "
        "user's IPTV provider playlist.",
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved: {OUT} ({len(prs.slides)} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
